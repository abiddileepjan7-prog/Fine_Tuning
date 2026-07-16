"""
=========================================================
PPTX Processing → Structured JSON + Timing Report
Library:
1. python-pptx
=========================================================
"""

import json
import os
import time
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

FILE = "Copy of Stocks Trading Business Plan.pptx"
OUTPUT_JSON = "pptx_structured_output.json"
TIMING_JSON = "pptx_timing_report.json"
IMAGE_DIR = "pptx_images"

os.makedirs(IMAGE_DIR, exist_ok=True)

# ==========================================================
# Unified output schema
# ==========================================================

result = {
    "source_file": FILE,
    "num_slides": 0,
    "slides": []
}

# ==========================================================
# Timing tracker
# ==========================================================

timings = {}


def time_it(label, func):
    """Runs func(), records elapsed time, returns func's result."""

    start = time.perf_counter()
    output = func()
    end = time.perf_counter()

    elapsed = round(end - start, 4)

    timings[label] = {"time_seconds": elapsed}

    print(f"\n[TIMER] {label} took {elapsed} seconds")

    return output


# ==========================================================
# 0. LOAD PRESENTATION (kept as its own timed step)
# ==========================================================

print("=" * 70)
print("PYTHON-PPTX")
print("=" * 70)


def run_load():

    prs = Presentation(FILE)

    result["num_slides"] = len(prs.slides)

    for slide_no, slide in enumerate(prs.slides, start=1):

        result["slides"].append({
            "slide_number": slide_no,
            "text": [],
            "tables": [],
            "charts": [],
            "images": []
        })

    return prs


prs = time_it("python_pptx_load", run_load)

timings["python_pptx_load"]["slides_found"] = result["num_slides"]

print("\nNumber of Slides :", result["num_slides"])

# ==========================================================
# 1. TEXT EXTRACTION
# ==========================================================

print("\n")
print("=" * 70)
print("1. TEXT EXTRACTION")
print("=" * 70)


def run_text_extraction():

    total_chars = 0

    for slide_no, slide in enumerate(prs.slides, start=1):

        slide_entry = result["slides"][slide_no - 1]

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:
                    slide_entry["text"].append(text)
                    total_chars += len(text)

    return total_chars


chars_text = time_it("python_pptx_text", run_text_extraction)

timings["python_pptx_text"]["chars_extracted"] = chars_text

print("Characters Extracted :", chars_text)

# ==========================================================
# 2. TABLE EXTRACTION
# ==========================================================

print("\n")
print("=" * 70)
print("2. TABLE EXTRACTION")
print("=" * 70)


def run_table_extraction():

    total_tables = 0

    for slide_no, slide in enumerate(prs.slides, start=1):

        slide_entry = result["slides"][slide_no - 1]

        for shape in slide.shapes:

            if shape.has_table:

                table = shape.table

                rows = []

                for row in table.rows:
                    rows.append([cell.text.strip() for cell in row.cells])

                slide_entry["tables"].append(rows)
                total_tables += 1

    return total_tables


tables_found = time_it("python_pptx_tables", run_table_extraction)

timings["python_pptx_tables"]["tables_found"] = tables_found

print("Tables Found :", tables_found)

# ==========================================================
# 3. CHART EXTRACTION
# ==========================================================

print("\n")
print("=" * 70)
print("3. CHART EXTRACTION")
print("=" * 70)


def run_chart_extraction():

    total_charts = 0

    for slide_no, slide in enumerate(prs.slides, start=1):

        slide_entry = result["slides"][slide_no - 1]

        for shape in slide.shapes:

            if shape.has_chart:

                chart = shape.chart

                chart_data = {
                    "type": str(chart.chart_type),
                    "categories": [],
                    "series": {}
                }

                try:
                    chart_data["categories"] = list(chart.plots[0].categories)

                    for series in chart.plots[0].series:
                        chart_data["series"][series.name] = list(series.values)

                except Exception as e:
                    chart_data["error"] = str(e)

                slide_entry["charts"].append(chart_data)
                total_charts += 1

    return total_charts


charts_found = time_it("python_pptx_charts", run_chart_extraction)

timings["python_pptx_charts"]["charts_found"] = charts_found

print("Charts Found :", charts_found)

# ==========================================================
# 4. IMAGE EXTRACTION
# ==========================================================

print("\n")
print("=" * 70)
print("4. IMAGE EXTRACTION")
print("=" * 70)


def run_image_extraction():

    total_images = 0

    for slide_no, slide in enumerate(prs.slides, start=1):

        slide_entry = result["slides"][slide_no - 1]

        image_count = 0

        for shape in slide.shapes:

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

                image_count += 1

                image = shape.image
                ext = image.ext

                filename = f"{IMAGE_DIR}/slide{slide_no}_image{image_count}.{ext}"

                with open(filename, "wb") as f:
                    f.write(image.blob)

                slide_entry["images"].append(filename)
                total_images += 1

    return total_images


images_found = time_it("python_pptx_images", run_image_extraction)

timings["python_pptx_images"]["images_found"] = images_found

print("Images Found :", images_found)

# ==========================================================
# SAVE FINAL STRUCTURED JSON
# ==========================================================

with open(OUTPUT_JSON, "w", encoding="utf-8") as f:
    json.dump(result, f, indent=2, ensure_ascii=False)

# ==========================================================
# TIMING REPORT — print to terminal + save to file
# ==========================================================

timing_report = {
    "source_file": FILE,
    "timings": timings
}

print("\n")
print("=" * 70)
print("TIMING REPORT")
print("=" * 70)

print(f"\n{'Step':<22}{'Time (s)'}")
print("-" * 35)

for step, data in timings.items():
    print(f"{step:<22}{data['time_seconds']}")

print("\nFull Timing Report (JSON)")
print(json.dumps(timing_report, indent=2, ensure_ascii=False))

with open(TIMING_JSON, "w", encoding="utf-8") as f:
    json.dump(timing_report, f, indent=2, ensure_ascii=False)

print(f"\nTiming report saved to {TIMING_JSON}")

print("\nDone.")